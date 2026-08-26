import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE  # Specifies the type of a chart
from pptx.dml.color import RGBColor


# fixed palette so bar colors are identical across every viewer (PowerPoint,
# Keynote, LibreOffice) instead of each one applying its own default theme
CHART_PALETTE = [
    RGBColor(0x1F, 0x77, 0xB4),
    RGBColor(0xFF, 0x7F, 0x0E),
    RGBColor(0x2C, 0xA0, 0x2C),
    RGBColor(0xD6, 0x27, 0x28),
    RGBColor(0x94, 0x67, 0xBD),
    RGBColor(0x8C, 0x56, 0x4B),
    RGBColor(0xE3, 0x77, 0xC2),
    RGBColor(0x7F, 0x7F, 0x7F),
    RGBColor(0xBC, 0xBD, 0x22),
    RGBColor(0x17, 0xBE, 0xCF),
]


class DeckBuilder:
    """Builds a PowerPoint deck slide-by-slide from Statista chat analysis data.

    Wraps a single python-pptx Presentation object. Each add_*_slide method
    appends one slide to it and returns the new slide. Call save() once all
    slides have been added to write the deck to disk.
    """

    def __init__(self):
        self.prs = Presentation()
        self.blank_layout = self.prs.slide_layouts[6]  # DRY principle

    def add_title_slide(self, title, subtitle=""):
        """Add the deck's opening slide: a large title and optional subtitle.

        Args:
            title: main title text.
            subtitle: optional subtitle text; skipped entirely if empty.

        Returns:
            The newly created slide.
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(3.9), Inches(8.4), Inches(0.6))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.word_wrap = True
            subtitle_frame.paragraphs[0].font.size = Pt(20)

        return slide

    def add_chart_slide(self, category, insight_text, metrics, source):
        """Add a slide with a title, an optional insight line, a bar chart
        built from `metrics`, and a source citation.

        Args:
            category: slide title text (the insight category name).
            insight_text: optional one-line summary shown under the title;
                skipped entirely if falsy.
            metrics: list of dicts, each expected to have a label under one of
                'action'/'statement'/'brand' and a numeric 'percentage'. Not a
                list, or containing non-dict items, is handled gracefully
                (falls back to a "No chart data available" placeholder rather
                than crashing) -- see EDGE_CASE.md.
            source: dict with citation fields (name, conductor,
                publication_date); if falsy, no citation is rendered.

        Returns:
            The newly created slide.
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        # --- title ---
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = category
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        # --- one-line insight, under the title ---
        if insight_text:
            insight_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.95), Inches(9), Inches(0.5))
            insight_frame = insight_box.text_frame
            insight_frame.text = insight_text
            insight_frame.word_wrap = True
            insight_frame.paragraphs[0].font.size = Pt(14)
            insight_frame.paragraphs[0].font.italic = True

        # --- extract labels + values from metrics, guarding against malformed input ---
        labels = []
        values = []
        # EDGE CASE: metrics could be None, a string, or anything else that's not
        # a list — skip the loop entirely instead of crashing on `for m in metrics`
        if isinstance(metrics, list):
            for m in metrics:
                # EDGE CASE: an item in the list could be a string/number instead
                # of a dict — skip just that item instead of crashing on m.get(...)
                if not isinstance(m, dict):
                    continue
                label = m.get('action') or m.get('statement') or m.get('brand') or 'Unknown'
                labels.append(label)
                values.append(m.get('percentage', 0))

        # EDGE CASE: metrics was empty/malformed and produced no usable data —
        # show a placeholder instead of handing an empty chart to python-pptx
        if labels and values:
            chart_data = CategoryChartData()
            chart_data.categories = labels
            chart_data.add_series('Percentage', values)

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.5),
                chart_data
            )

            # explicitly color each bar so the palette is baked into the file
            # itself (same colors in pptx and in the PDF conversion), instead
            # of relying on each viewer's own default theme
            plot = chart_frame.chart.plots[0]
            plot.vary_by_categories = True
            series = plot.series[0]
            for i, point in enumerate(series.points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]

            value_axis = chart_frame.chart.value_axis
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = "Percentage (%)"

            category_axis = chart_frame.chart.category_axis
            category_axis.tick_labels.font.size = Pt(10)
        else:
            placeholder_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.6), Inches(8.4), Inches(1))
            placeholder_frame = placeholder_box.text_frame
            placeholder_frame.text = "No chart data available"
            placeholder_frame.word_wrap = True
            placeholder_frame.paragraphs[0].font.size = Pt(16)
            placeholder_frame.paragraphs[0].font.italic = True

        # --- source citation, bottom of slide ---
        if source:
            citation_text = (
                f"Source: {source.get('name', 'Unknown')} — "
                f"{source.get('conductor', '')}, {source.get('publication_date', '')}"
            )
            citation_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
            citation_frame = citation_box.text_frame
            citation_frame.text = citation_text
            citation_frame.word_wrap = True
            citation_frame.paragraphs[0].font.size = Pt(10)
            citation_frame.paragraphs[0].font.italic = True

        return slide

    def add_summary_slide(self, summary_text, caveat=None):
        """Add a closing slide: `summary_text` split into bullet points, plus
        an optional caveat note about data quality.

        Args:
            summary_text: a block of prose; split on paragraph breaks then
                sentence boundaries into one bullet per sentence (a simple
                heuristic -- see EDGE_CASE.md for its known limitations).
            caveat: optional string rendered as a small italic note at the
                bottom of the slide; skipped entirely if falsy.

        Returns:
            The newly created slide.
        """
        slide = self.prs.slides.add_slide(self.blank_layout)

        # --- title ---
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Summary & Key Takeaways"
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        # --- summary, split into bullet-style sentences ---
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.2))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True

        # split on paragraph breaks first (colon + \n\n isn't caught by '. ' alone),
        # then split each paragraph into sentences
        paragraphs = summary_text.split('\n\n')
        sentences = []
        for para in paragraphs:
            para = para.replace('\n', ' ')
            for s in para.split('. '):
                s = s.strip()
                if s:
                    sentences.append(s)

        for i, sentence in enumerate(sentences):
            if not sentence.endswith('.'):
                sentence += '.'
            if i == 0:
                paragraph = body_frame.paragraphs[0]
            else:
                paragraph = body_frame.add_paragraph()
            paragraph.text = f"• {sentence}"
            paragraph.font.size = Pt(14)
            paragraph.space_after = Pt(10)

        # --- caveat, if provided ---
        if caveat:
            caveat_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.7))
            caveat_frame = caveat_box.text_frame
            caveat_frame.text = f"Note: {caveat}"
            caveat_frame.word_wrap = True
            caveat_frame.paragraphs[0].font.size = Pt(11)
            caveat_frame.paragraphs[0].font.italic = True

        return slide

    def save(self, path):
        """Write the deck to disk at `path`, creating parent directories if
        they don't already exist.

        Args:
            path: destination file path, e.g. 'example_output/statista_ppt.pptx'.
        """
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)
